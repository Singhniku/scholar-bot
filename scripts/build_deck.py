"""
Generate a Scholar-Bot demo presentation (.pptx).
Run:  python scripts/build_deck.py
Output: docs/Scholar-Bot-Demo.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Theme ─────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
BLUE   = RGBColor(0x2D, 0x6A, 0x9F)
LIGHT  = RGBColor(0xF0, 0xF6, 0xFF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x55, 0x55, 0x55)
DARK   = RGBColor(0x14, 0x14, 0x14)
GREEN  = RGBColor(0x15, 0x57, 0x24)
ORANGE = RGBColor(0x85, 0x64, 0x04)
RED    = RGBColor(0x72, 0x1C, 0x24)
ACCENT = RGBColor(0xFF, 0x6B, 0x35)

# 16:9 widescreen
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "screenshots"


# ── Helpers ───────────────────────────────────────────────────────────────────
def add_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def line(shape, rgb=None, width=None):
    if rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb
        if width: shape.line.width = width


def text(shape, content, *, size=14, bold=False, color=DARK,
         align=PP_ALIGN.LEFT, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = content
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def textbox(slide, left, top, width, height, content, **kw):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    text(box, content, **kw)
    return box


def add_box(slide, l, t, w, h, content, *, fillc=BLUE, textc=WHITE,
            size=14, bold=True, align=PP_ALIGN.CENTER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    fill(s, fillc); line(s, NAVY, Pt(0.75))
    s.text_frame.margin_left = s.text_frame.margin_right = Inches(0.05)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(s, content, size=size, bold=bold, color=textc, align=align)
    return s


def title_bar(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, 0, prs.slide_width, Inches(1.0))
    fill(bar, NAVY); line(bar)
    bar.text_frame.margin_left = Inches(0.5)
    bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(bar, title, size=28, bold=True, color=WHITE)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.55),
                                       Inches(12), Inches(0.4))
        text(sb, subtitle, size=13, color=LIGHT)


def add_arrow(slide, x1, y1, x2, y2, color=BLUE, width=Pt(2.5)):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x1), Inches(y1),
                                       Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = width
    # Add arrow head via XML
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = conn.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle"); tail.set("w", "med"); tail.set("h", "med")


# ════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                        prs.slide_width, prs.slide_height)
fill(bg, NAVY); line(bg)

# Logo block
logo = s.shapes.add_shape(MSO_SHAPE.OVAL,
                           Inches(5.7), Inches(1.4), Inches(2), Inches(2))
fill(logo, ACCENT); line(logo)
logo.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
text(logo, "🎓", size=72, color=WHITE, align=PP_ALIGN.CENTER)

textbox(s, 1, 3.6, 11.3, 1, "Scholar-Bot",
        size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, 1, 4.6, 11.3, 0.7,
        "AI-Powered Resume Optimizer · LinkedIn Job Matcher · Auto-Apply Agent",
        size=20, color=LIGHT, align=PP_ALIGN.CENTER)
textbox(s, 1, 6.4, 11.3, 0.5,
        "github.com/Singhniku/scholar-bot",
        size=14, color=LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 2 — The Problem
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "The Problem",
          "Job seekers waste 10+ hours per week customising resumes for ATS")

problems = [
    ("75%", "of resumes are auto-rejected by ATS before any human reads them"),
    ("3-5 min", "average time to manually customise one resume per job"),
    ("100s", "of jobs posted daily — manual review is unsustainable"),
    ("Generic", "applications get lower interview rates than tailored ones"),
]
for i, (stat, line_) in enumerate(problems):
    col = i % 2; row = i // 2
    l = 0.7 + col * 6.4
    t = 1.6 + row * 2.6

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(l), Inches(t), Inches(5.9), Inches(2.3))
    fill(box, LIGHT); line(box, BLUE, Pt(2))

    textbox(s, l + 0.2, t + 0.2, 5.5, 0.9, stat,
            size=44, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    textbox(s, l + 0.2, t + 1.2, 5.5, 1.0, line_,
            size=15, color=DARK, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 3 — The Solution
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "The Solution",
          "End-to-end AI automation: parse → match → rewrite → apply")

steps = [
    ("📄", "Parse",     "Read PDF / JPG / PNG / SVG resumes via pdfplumber + OCR"),
    ("🧠", "Extract",   "AI extracts skills, experience, education as structured JSON"),
    ("🔍", "Search",    "Scrape LinkedIn for jobs matching your title + filters"),
    ("📊", "Score",     "Weighted match algorithm against required + preferred skills"),
    ("✨", "Rewrite",   "AI tailors bullet points to mirror each job's exact keywords"),
    ("🚀", "Auto-Apply","Selenium fills Easy Apply forms — you approve, it submits"),
]
for i, (icon, title_, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    l = 0.5 + col * 4.3
    t = 1.5 + row * 2.85

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(l), Inches(t), Inches(4.0), Inches(2.6))
    fill(box, WHITE); line(box, BLUE, Pt(1.5))
    textbox(s, l + 0.1, t + 0.1, 1, 0.7, icon,
            size=32, align=PP_ALIGN.CENTER)
    textbox(s, l + 0.1, t + 0.85, 3.8, 0.5, title_,
            size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    textbox(s, l + 0.2, t + 1.4, 3.6, 1.1, desc,
            size=12, color=DARK, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 4 — Architecture
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "Architecture",
          "Layered design — UI · Pipeline · Services · External integrations")

# UI layer (top, full width)
add_box(s, 0.5, 1.3, 12.3, 0.8,
        "Streamlit UI  ·  5 tabs: Upload · Job Matches · Optimised Resumes · Auto Apply · Report",
        fillc=NAVY, size=15)

# Pipeline layer
add_box(s, 0.5, 2.3, 12.3, 0.8,
        "Pipeline (src/pipeline.py)  ·  parse → extract → fetch → match → filter → upgrade",
        fillc=BLUE, size=14)

# Service boxes (4 columns)
services = [
    ("Resume Parser",      "pdfplumber\npytesseract\ncairosvg"),
    ("Skills Extractor",   "AI Client\n(Gemini / Claude)\n+ Keyword fallback"),
    ("LinkedIn Scraper",   "requests\nBeautifulSoup\nlxml"),
    ("ATS Optimiser",      "AI Client\n+ ResumeGenerator\n(reportlab PDF)"),
]
for i, (name, lib) in enumerate(services):
    l = 0.5 + i * 3.1
    add_box(s, l, 3.3, 2.95, 1.5, name, fillc=LIGHT, textc=DARK, size=13)
    textbox(s, l + 0.05, 3.85, 2.85, 1.0, lib,
            size=10, color=GREY, align=PP_ALIGN.CENTER)

# External
add_box(s, 0.5, 5.0, 3.95, 0.9, "Google Gemini\n(free tier, default)",
        fillc=GREEN, size=12)
add_box(s, 4.7, 5.0, 3.95, 0.9, "Anthropic Claude\n(paid, optional)",
        fillc=BLUE, size=12)
add_box(s, 8.9, 5.0, 3.9, 0.9, "LinkedIn Public Jobs\n+ Easy Apply (Selenium)",
        fillc=NAVY, size=12)

# Storage
add_box(s, 0.5, 6.1, 12.3, 0.8,
        "Output: PDF · Markdown · JSON · CSV  ·  Cached in /output, /uploads",
        fillc=GREY, size=12)


# ════════════════════════════════════════════════════════════════════════════
# Slide 5 — Pipeline Flow Diagram
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "Pipeline Flow",
          "Six pure functions in src/pipeline.py — same code path for UI, CLI, tests")

stages = [
    ("Resume\nFile", 0.5, 2.3, 1.5, 1.2, ACCENT),
    ("parse_resume()\n→ raw text", 2.3, 2.3, 2.0, 1.2, BLUE),
    ("extract_resume_skills()\n→ resume_data", 4.6, 2.3, 2.2, 1.2, BLUE),
    ("fetch_jobs()\n→ list[job]", 7.1, 2.3, 2.0, 1.2, BLUE),
    ("match_resume_to_jobs()\n→ list[ScoredJob]", 9.4, 2.3, 2.4, 1.2, BLUE),
]
for label, l, t, w, h, c in stages:
    add_box(s, l, t, w, h, label, fillc=c, size=11)

# Arrows between row 1
for x in (2.0, 4.3, 6.8, 9.1):
    add_arrow(s, x, 2.9, x + 0.3, 2.9)

# Row 2
stages2 = [
    ("filter_by_match()\n→ above / below", 0.5, 4.2, 2.5, 1.2, BLUE),
    ("upgrade_cv()\n→ optimised resume", 3.3, 4.2, 2.5, 1.2, ACCENT),
    ("ResumeGenerator\n→ PDF / Markdown", 6.1, 4.2, 2.5, 1.2, BLUE),
    ("AutoApply\n→ Easy Apply form", 8.9, 4.2, 2.5, 1.2, GREEN),
]
for label, l, t, w, h, c in stages2:
    add_box(s, l, t, w, h, label, fillc=c, size=11)

for x in (3.0, 5.8, 8.6):
    add_arrow(s, x, 4.8, x + 0.3, 4.8)

# Down arrow
add_arrow(s, 11.8, 3.5, 1.7, 4.2)

# Footer
textbox(s, 0.5, 6.0, 12.3, 0.6,
        "Each step is independent · Falls back gracefully when AI quota runs out",
        size=14, color=GREY, align=PP_ALIGN.CENTER)
textbox(s, 0.5, 6.6, 12.3, 0.6,
        "61 unit + integration tests · 100% pipeline coverage",
        size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 6 — Features
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "Features",
          "Battle-tested features for real-world job hunting")

features = [
    ("📄", "Multi-format parsing", "PDF (native + OCR fallback), JPG, PNG, SVG"),
    ("🤖", "Dual AI provider",     "Google Gemini (free) or Anthropic Claude (paid)"),
    ("🛡️", "Keyword fallback",    "Fully usable when AI quota is exhausted"),
    ("🔍", "Title + filter search","Job title, location, days, min match %"),
    ("📊", "ATS audit",            "7-rule keyword score + AI bullet rewrites"),
    ("✨", "ATS-safe rewrites",    "Mirror exact JD keywords without fabrication"),
    ("📈", "Before/after ATS",     "See exactly how much your score improved"),
    ("🚀", "Human-gate auto-apply","Bot fills, you approve, it submits"),
]
for i, (icon, title_, desc) in enumerate(features):
    col = i % 2; row = i // 2
    l = 0.5 + col * 6.4
    t = 1.4 + row * 1.4

    icon_box = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(l), Inches(t), Inches(0.9), Inches(0.9))
    fill(icon_box, LIGHT); line(icon_box, BLUE, Pt(1.5))
    icon_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(icon_box, icon, size=22, align=PP_ALIGN.CENTER)

    textbox(s, l + 1.1, t, 5.2, 0.5, title_,
            size=16, bold=True, color=NAVY)
    textbox(s, l + 1.1, t + 0.5, 5.2, 0.7, desc,
            size=12, color=DARK)


# ════════════════════════════════════════════════════════════════════════════
# Slides 7-11 — Screenshots
# ════════════════════════════════════════════════════════════════════════════
def screenshot_slide(title, subtitle, image_path, caption):
    s = add_blank()
    title_bar(s, title, subtitle)
    if image_path.exists():
        # Centre-fit the image in the lower 6 inches of the slide
        s.shapes.add_picture(str(image_path),
                              Inches(1.5), Inches(1.3),
                              Inches(10.3), Inches(5.5))
    textbox(s, 0.5, 6.85, 12.3, 0.5, caption,
            size=12, color=GREY, align=PP_ALIGN.CENTER)


screenshot_slide(
    "Tab 1: Upload Resume",
    "Drag & drop · Sidebar configures provider, search, auto-apply",
    SHOTS / "01-upload-resume.png",
    "Skills, experience, education extracted instantly. ATS audit shows 90/100 with 7-rule breakdown."
)

screenshot_slide(
    "Tab 2: Job Matches",
    "Ranked LinkedIn results with match score, expandable details, per-job actions",
    SHOTS / "02-job-matches.png",
    "Each job shows title, company, location, posting date · 'Upgrade CV' button per job"
)

screenshot_slide(
    "Tab 3: Optimised Resume",
    "Full preview with before/after ATS score and downloadable PDF/Markdown",
    SHOTS / "03-optimized-resumes.png",
    "AI rewrites every bullet to mirror job keywords · downloads PDF + Markdown"
)

screenshot_slide(
    "Tab 4: Auto Apply",
    "Selenium fills LinkedIn Easy Apply forms — pauses for your approval before submit",
    SHOTS / "04-auto-apply.png",
    "Human-in-the-loop gate · threading.Event blocks the bot until you click submit"
)

screenshot_slide(
    "Tab 5: Report",
    "Ranked job table with match scores, distribution chart, CSV export",
    SHOTS / "05-report.png",
    "Full audit trail of every job seen, scored, queued, applied"
)


# ════════════════════════════════════════════════════════════════════════════
# Slide 12 — AI Models & Libraries
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "AI Models & Libraries",
          "Production-grade open-source stack")

# Two columns
textbox(s, 0.5, 1.4, 6, 0.5, "🤖 AI Models",
        size=20, bold=True, color=NAVY)

ai_models = [
    ("Google Gemini 2.5 Flash",  "Default · Free tier · 1500 req/day",  GREEN),
    ("Gemini 2.0 Flash",         "Auto-fallback when 2.5 is rate-limited", BLUE),
    ("Gemini 2.0 Flash-Lite",    "Final fallback in the chain",         BLUE),
    ("Anthropic Claude Sonnet 4.6","Optional · Paid · w/ prompt caching", NAVY),
]
for i, (name, desc, c) in enumerate(ai_models):
    t = 2.0 + i * 0.85
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(t), Inches(6), Inches(0.75))
    fill(box, LIGHT); line(box, c, Pt(1.5))
    textbox(s, 0.7, t + 0.05, 5.6, 0.35, name,
            size=14, bold=True, color=c)
    textbox(s, 0.7, t + 0.4, 5.6, 0.35, desc,
            size=11, color=GREY)

textbox(s, 6.8, 1.4, 6, 0.5, "📚 Key Libraries",
        size=20, bold=True, color=NAVY)

libs = [
    ("streamlit",        "1.57   · 5-tab UI"),
    ("google-genai",     "0.8    · Gemini SDK"),
    ("anthropic",        "0.97   · Claude SDK"),
    ("pdfplumber",       "0.11   · PDF text + tables"),
    ("pytesseract",      "0.3    · OCR for images"),
    ("cairosvg",         "2.9    · SVG → PNG"),
    ("beautifulsoup4",   "4.14   · LinkedIn HTML parse"),
    ("selenium",         "4.43   · Easy Apply automation"),
    ("reportlab",        "4.5    · ATS-safe PDF gen"),
]
for i, (lib, desc) in enumerate(libs):
    t = 2.0 + i * 0.5
    textbox(s, 6.9, t, 1.8, 0.4, lib,
            size=12, bold=True, color=NAVY, font="Consolas")
    textbox(s, 8.7, t, 4.2, 0.4, desc,
            size=11, color=DARK)


# ════════════════════════════════════════════════════════════════════════════
# Slide 13 — ATS Scoring Logic
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "ATS Scoring Logic",
          "Transparent, explainable 7-rule audit · 0-100")

rules = [
    ("Contact info present (name + email + phone)", "20"),
    ("Professional summary (≥30 chars)",            "15"),
    ("Skills section present",                       "15"),
    ("Work experience with bullet points",           "20"),
    ("Education section present",                    "10"),
    ("≥10 technical keywords detected",              "10"),
    ("Word count between 400-1200",                  "10"),
]
# Header
add_box(s, 1, 1.4, 9, 0.6, "Check", fillc=NAVY, size=14)
add_box(s, 10, 1.4, 2.3, 0.6, "Points", fillc=NAVY, size=14)

for i, (rule, pts) in enumerate(rules):
    t = 2.05 + i * 0.55
    bg = LIGHT if i % 2 == 0 else WHITE
    rect1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(1), Inches(t), Inches(9), Inches(0.5))
    fill(rect1, bg); line(rect1, BLUE, Pt(0.5))
    textbox(s, 1.2, t + 0.05, 8.6, 0.4, rule, size=12, color=DARK)

    rect2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(10), Inches(t), Inches(2.3), Inches(0.5))
    fill(rect2, bg); line(rect2, BLUE, Pt(0.5))
    textbox(s, 10.2, t + 0.05, 1.9, 0.4, pts,
            size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# AI augmentation
add_box(s, 1, 6.0, 11.3, 0.9,
        "🤖 With AI:  bullet rewrites · keyword gap analysis · format issue detection",
        fillc=GREEN, size=14)


# ════════════════════════════════════════════════════════════════════════════
# Slide 14 — Benefits
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "Benefits",
          "Hours saved · Higher response rate · Zero fabrication")

benefits = [
    ("⏱️", "10× faster",        "Customise 10 resumes in the time it takes to do 1 manually"),
    ("🎯", "ATS-optimised",     "Single-column PDFs with mirrored keywords pass automated filters"),
    ("🆓", "Free forever",       "Gemini's free tier handles 1500 req/day — no credit card needed"),
    ("🤝", "Human-in-the-loop","Auto-Apply pauses for approval before every single submission"),
    ("🔒", "Privacy-first",     "Resume processed locally · API keys stored in your .env only"),
    ("📈", "Measurable",        "Before/after ATS score + match analytics on every job"),
]
for i, (icon, title_, desc) in enumerate(benefits):
    col = i % 3; row = i // 3
    l = 0.5 + col * 4.3
    t = 1.4 + row * 2.95

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(l), Inches(t), Inches(4.1), Inches(2.7))
    fill(box, WHITE); line(box, ACCENT, Pt(2))

    textbox(s, l + 0.2, t + 0.3, 3.7, 0.7, icon,
            size=44, align=PP_ALIGN.CENTER)
    textbox(s, l + 0.2, t + 1.1, 3.7, 0.5, title_,
            size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    textbox(s, l + 0.2, t + 1.7, 3.7, 1.0, desc,
            size=12, color=DARK, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 15 — Tech Stack Summary
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "Tech Stack at a Glance",
          "Chosen for: free-tier-friendly, production-ready, easy to deploy")

layers = [
    ("Frontend",    "Streamlit (Python web UI)",                          NAVY),
    ("AI / LLM",    "google-genai · anthropic SDK · prompt caching",       BLUE),
    ("Parsing",     "pdfplumber · pytesseract · cairosvg · Pillow",       BLUE),
    ("Scraping",    "requests · beautifulsoup4 · lxml · python-dateutil", BLUE),
    ("Automation",  "selenium · webdriver-manager (Chrome)",              BLUE),
    ("Output",      "reportlab (PDF) · pandas (CSV)",                      BLUE),
    ("Testing",     "pytest · 61 unit + integration tests",                GREEN),
    ("Setup",       "setup.sh · 3-command install",                        ORANGE),
]
for i, (layer, tech, c) in enumerate(layers):
    t = 1.4 + i * 0.65
    add_box(s, 1, t, 2.5, 0.55, layer, fillc=c, size=12)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(3.7), Inches(t), Inches(8.5), Inches(0.55))
    fill(box, WHITE); line(box, c, Pt(1))
    textbox(s, 3.8, t + 0.05, 8.3, 0.45, tech,
            size=12, color=DARK)


# ════════════════════════════════════════════════════════════════════════════
# Slide 16 — Quality / Tests
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "Quality & Testing",
          "61 tests · zero-network mocked unit tests · live integration tests")

stats = [
    ("61",     "Tests passing", GREEN),
    ("9",      "Test categories",  BLUE),
    ("0",      "Failures",      GREEN),
    ("100%",   "Pipeline coverage", BLUE),
]
for i, (num, label, c) in enumerate(stats):
    l = 0.5 + i * 3.2
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(l), Inches(1.4), Inches(3), Inches(2.0))
    fill(box, WHITE); line(box, c, Pt(2))
    textbox(s, l + 0.2, 1.55, 2.6, 1.2, num,
            size=64, bold=True, color=c, align=PP_ALIGN.CENTER)
    textbox(s, l + 0.2, 2.7, 2.6, 0.5, label,
            size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

textbox(s, 0.5, 3.8, 12.3, 0.4, "Test categories:",
        size=16, bold=True, color=NAVY)
categories = [
    "✓ Fallback resume parser — positive cases",
    "✓ Edge cases — empty, no dates, single role",
    "✓ Overlapping date-range merging",
    "✓ Normalizer type coercion (\"5+\" → 5, \"5 yrs\" → 5)",
    "✓ LinkedIn query construction",
    "✓ Score match — perfect, zero, partial",
    "✓ Job description extraction",
    "✓ PDF + Markdown generation",
    "✓ Live LinkedIn search — title vs skills",
    "✓ Pipeline orchestration — 8 tests for src/pipeline.py",
]
for i, cat in enumerate(categories):
    col = i % 2; row = i // 2
    textbox(s, 0.7 + col * 6.2, 4.3 + row * 0.4, 6, 0.4, cat,
            size=12, color=DARK)


# ════════════════════════════════════════════════════════════════════════════
# Slide 17 — Setup
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "Quick Start",
          "From clone to running app in under 2 minutes")

# Code block
code = (
    "$ git clone https://github.com/Singhniku/scholar-bot.git\n"
    "$ cd scholar-bot\n"
    "$ ./setup.sh                              ← creates venv, installs deps,\n"
    "                                            tesseract, .env, folders\n"
    "$ source .venv/bin/activate\n"
    "$ streamlit run app.py                    ← open localhost:8501"
)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.7), Inches(1.4), Inches(11.9), Inches(2.8))
fill(box, DARK); line(box, NAVY, Pt(1))
tf = box.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = code
run.font.name = "Consolas"
run.font.size = Pt(15)
run.font.color.rgb = RGBColor(0xC8, 0xE6, 0xC9)

# 3 steps
steps3 = [
    ("1.", "Clone",   "git clone the repo"),
    ("2.", "Setup",   "./setup.sh — one command"),
    ("3.", "Run",     "streamlit run app.py"),
]
for i, (num, title_, desc) in enumerate(steps3):
    l = 0.7 + i * 4.05
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(l), Inches(4.5), Inches(3.85), Inches(2.4))
    fill(box, LIGHT); line(box, BLUE, Pt(2))
    textbox(s, l + 0.1, 4.6, 3.65, 0.7, num,
            size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    textbox(s, l + 0.1, 5.55, 3.65, 0.5, title_,
            size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    textbox(s, l + 0.1, 6.15, 3.65, 0.7, desc,
            size=12, color=DARK, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 18 — Roadmap / Q&A
# ════════════════════════════════════════════════════════════════════════════
s = add_blank()
title_bar(s, "Roadmap & Future Work",
          "Where Scholar-Bot is heading next")

next_up = [
    "🌐  Indeed · Glassdoor · Naukri scrapers",
    "✉️   AI-generated cover letter per job",
    "📊  Application tracking dashboard",
    "📝  .docx resume input",
    "🐳  Docker image for one-command deploy",
    "🌍  Multi-language resume support",
]
for i, item in enumerate(next_up):
    col = i % 2; row = i // 2
    l = 0.7 + col * 6.2
    t = 1.5 + row * 0.9
    textbox(s, l, t, 6, 0.7, item,
            size=18, color=DARK)

# Closing
add_box(s, 0.7, 5.4, 11.9, 1.6,
        "🎓  github.com/Singhniku/scholar-bot   ·   Try it · Fork it · Star it",
        fillc=NAVY, size=20)


# ── Save ──────────────────────────────────────────────────────────────────────
out = ROOT / "docs" / "Scholar-Bot-Demo.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"✓ Saved: {out}  ({out.stat().st_size // 1024} KB · {len(prs.slides)} slides)")
